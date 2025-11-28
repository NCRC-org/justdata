#!/usr/bin/env python3
"""
PowerPoint export functionality for DataExplorer reports.
Creates PowerPoint files in 16:9 format with NCRC branding.
"""

import os
from typing import Dict, Any, List, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# NCRC Brand Colors
NCRC_PURPLE = RGBColor(85, 45, 135)  # #552d87
NCRC_BLUE = RGBColor(3, 78, 160)  # #034ea0
NCRC_LIGHT_BLUE = RGBColor(47, 173, 227)  # #2fade3
NCRC_RED = RGBColor(232, 46, 46)  # #e82e2e
NCRC_WHITE = RGBColor(255, 255, 255)
NCRC_BLACK = RGBColor(0, 0, 0)
NCRC_GRAY = RGBColor(129, 131, 144)  # #818390


def save_dataexplorer_powerpoint_report(
    tables: Dict[str, List[Dict[str, Any]]],
    filters: Dict[str, Any],
    output_path: str,
    report_type: str = 'area',  # 'area' or 'lender'
    metadata: Optional[Dict[str, Any]] = None
):
    """
    Save the DataExplorer report data to a PowerPoint file in 16:9 format.
    
    Args:
        tables: Dictionary of table names to table data (list of dicts)
        filters: Dictionary of filters applied to the report
        output_path: Path to save the PowerPoint file
        report_type: Type of report ('area' or 'lender')
        metadata: Optional metadata about the report
    """
    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    
    # Create presentation in 16:9 format
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio
    
    # Get metadata
    data_type = filters.get('dataType', 'hmda')
    years = filters.get('years', [])
    years_str = f"{min(years)}-{max(years)}" if years else "N/A"
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"DataExplorer {report_type.title()} Analysis"
    subtitle.text = f"{data_type.upper()} Data - {years_str}"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.color.rgb = NCRC_PURPLE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.color.rgb = NCRC_GRAY
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Methods/Definitions slide
    methods_slide_layout = prs.slide_layouts[5]  # Blank layout
    methods_slide = prs.slides.add_slide(methods_slide_layout)
    
    # Add title
    title_shape = methods_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.text = "Methods and Definitions"
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = NCRC_PURPLE
    
    # Add content
    content_y = Inches(1.2)
    content_shape = methods_slide.shapes.add_textbox(Inches(0.5), content_y, Inches(9), Inches(4))
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True
    
    methods_text = []
    methods_text.append(f"Report Type: {report_type.title()} Analysis")
    methods_text.append(f"Data Type: {data_type.upper()}")
    methods_text.append(f"Years: {years_str}")
    if metadata:
        if 'geography' in metadata:
            methods_text.append(f"Geography: {metadata['geography']}")
        if 'subject_lender' in metadata:
            methods_text.append(f"Subject Lender: {metadata['subject_lender']}")
    methods_text.append("")
    methods_text.append("Data Sources:")
    if data_type == 'hmda':
        methods_text.append("• Home Mortgage Disclosure Act (HMDA) data")
    elif data_type == 'sb':
        methods_text.append("• CRA Small Business Lending data")
    elif data_type == 'branches':
        methods_text.append("• FDIC Summary of Deposits (SOD) branch data")
    methods_text.append("• U.S. Census Bureau ACS 5-Year Estimates")
    methods_text.append("")
    methods_text.append("For detailed definitions and filter information, see the Excel export Methods sheet.")
    
    for i, line in enumerate(methods_text):
        if i > 0:
            content_frame.add_paragraph()
        para = content_frame.paragraphs[i]
        para.text = line
        para.font.size = Pt(14)
        para.font.color.rgb = NCRC_BLACK
        if line.endswith(":") or (line.startswith("Report") or line.startswith("Data") or line.startswith("Years")):
            para.font.bold = True
    
    # Create a slide for each table
    for table_name, table_data in tables.items():
        if table_data:
            # Use blank layout
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            title_frame = title_shape.text_frame
            title_frame.text = table_name
            title_para = title_frame.paragraphs[0]
            title_para.font.bold = True
            title_para.font.size = Pt(28)
            title_para.font.color.rgb = NCRC_PURPLE
            
            # Create table
            # Determine table dimensions
            num_rows = len(table_data) + 1  # +1 for header
            if num_rows == 1:
                continue  # Skip empty tables
            
            # Get all unique keys from all rows
            all_keys = set()
            for row in table_data:
                all_keys.update(row.keys())
            num_cols = len(all_keys)
            
            if num_cols == 0:
                continue
            
            # Limit columns to fit on slide
            max_cols = 8
            if num_cols > max_cols:
                # Take first max_cols columns
                all_keys = list(all_keys)[:max_cols]
                num_cols = max_cols
            
            # Calculate table dimensions
            table_left = Inches(0.5)
            table_top = Inches(1.2)
            table_width = Inches(9)
            table_height = min(Inches(0.5) * num_rows, Inches(3.5))
            
            # Create table
            table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
            table = table_shape.table
            
            # Set column widths
            col_width = table_width / num_cols
            for col_idx in range(num_cols):
                table.columns[col_idx].width = int(col_width)
            
            # Fill header row
            header_keys = list(all_keys)
            for col_idx, key in enumerate(header_keys):
                cell = table.cell(0, col_idx)
                cell.text = str(key)
                # Style header
                cell.fill.solid()
                cell.fill.fore_color.rgb = NCRC_PURPLE
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = NCRC_WHITE
                para.font.size = Pt(10)
                para.alignment = PP_ALIGN.CENTER
            
            # Fill data rows
            for row_idx, row_data in enumerate(table_data, start=1):
                for col_idx, key in enumerate(header_keys):
                    cell = table.cell(row_idx, col_idx)
                    value = row_data.get(key, '')
                    cell.text = str(value) if value is not None else ''
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(9)
                    para.alignment = PP_ALIGN.CENTER
                    
                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(250, 250, 250)  # Light gray
            
            # Add border to table
            for row in table.rows:
                for cell in row.cells:
                    # Set border
                    cell.fill.fore_color.rgb = None if cell.fill.fore_color.rgb == RGBColor(250, 250, 250) else None
    
    # Save presentation
    prs.save(output_path)

